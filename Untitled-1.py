# =============================================================
# SignBridge – Final Presentation Generator
# Mapped to Origin Template with Dark Space / Magenta Theme
# =============================================================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette (Dark Space / Magenta Theme) ──
DARK_BG     = RGBColor(10, 5, 15)         # Deep space black/dark purple
WHITE       = RGBColor(255, 255, 255)     # Main text
TEAL        = RGBColor(160, 32, 112)      # Deep Magenta / Purple
CORAL       = RGBColor(220, 100, 150)     # Vibrant Pink / Light Magenta
LIGHT_BG    = RGBColor(40, 15, 45)        # Dark magenta-purple for cards
DARK_GRAY   = RGBColor(220, 220, 220)     # Light gray for standard text
ACCENT_BLUE = RGBColor(0, 150, 255)       # Bright blue accent
GOLD        = RGBColor(230, 190, 255)     # Pale pink/purple

def add_background_shape(slide, prs, color):
    """Add a full-slide colored rectangle as background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape

def add_accent_bar(slide, left, top, width, height, color):
    """Add a colored accent bar/stripe."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Times New Roman", line_spacing=1.2):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox

def add_bullet_points(slide, left, top, width, height, items, font_size=18):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "Times New Roman"
        p.space_after = Pt(10)
        p.level = 0
    return txBox

# =============================================================
# CREATE PRESENTATION
# =============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_WIDTH = prs.slide_width

# Helper for standard slide header
def setup_standard_slide(slide, title_text):
    add_background_shape(slide, prs, DARK_BG)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.1), TEAL)
    add_accent_bar(slide, Inches(0), Inches(7.4), SLIDE_WIDTH, Inches(0.1), CORAL)
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8),
                 title_text, font_size=36, font_color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(0.05), CORAL)

# -------------------------------------------------------------
# SLIDE 1: Problem Statement
# -------------------------------------------------------------
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
setup_standard_slide(slide1, "Problem Statement")
bullets1 = [
    "Millions are silenced by communication barriers.",
    "Over 70 million people worldwide rely on sign language.",
    "Less than 1% of the general public understands sign language.",
    "90% of deaf children are born to hearing parents.",
    "According to the WHO, 466 million people have disabling hearing loss."
]
add_bullet_points(slide1, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4), bullets1, font_size=22)

# -------------------------------------------------------------
# SLIDE 2: Identifying Key Challenges We Face
# -------------------------------------------------------------
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
setup_standard_slide(slide2, "Identifying Key Challenges We Face")
bullets2 = [
    "🏥 Hospitals: Deaf patients cannot properly explain their symptoms to doctors.",
    "🏫 Schools: Students miss out on their education when interpreters are not available.",
    "🏢 Workplaces: Career growth is severely limited due to communication gaps.",
    "🏛️ Public Services: Accessing basic government services becomes a struggle."
]
add_bullet_points(slide2, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4), bullets2, font_size=22)

# -------------------------------------------------------------
# SLIDE 3: Proposed Solution
# -------------------------------------------------------------
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
setup_standard_slide(slide3, "Proposed Solution")
bullets3 = [
    "SignBridge: An AI-powered, real-time sign language translator.",
    "Utilizes computer vision and machine learning to instantly capture and recognize hand gestures.",
    "Bridges the communication gap by converting recognized signs into readable on-screen text and natural audible speech."
]
add_bullet_points(slide3, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4), bullets3, font_size=24)

# -------------------------------------------------------------
# SLIDE 4: A Path to Innovation
# -------------------------------------------------------------
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
setup_standard_slide(slide4, "A Path to Innovation")
bullets4 = [
    "Step 1: A user shows a hand sign in front of a mobile or webcam camera.",
    "Step 2: Computer vision detects the hand and tracks 21 specific finger landmarks.",
    "Step 3: A machine learning model compares the detected gesture against a trained sign language dataset.",
    "Step 4: The system predicts the corresponding word.",
    "Step 5: The predicted word is immediately displayed as text and spoken out loud via text-to-speech."
]
add_bullet_points(slide4, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4), bullets4, font_size=20)

# -------------------------------------------------------------
# SLIDE 5: System Architecture
# -------------------------------------------------------------
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
setup_standard_slide(slide5, "System Architecture")
bullets5 = [
    "The architecture follows a real-time data flow pipeline.",
    "Begins with camera input capturing a video feed, sent to MediaPipe for hand detection and extraction of 21 3D hand landmarks per frame.",
    "The data flows into an ML model that classifies the gesture.",
    "Outputs the predicted word onto the screen and speaks it using pyttsx3 or gTTS.",
    "Processing runs efficiently at 30+ FPS to ensure real-time performance."
]
add_bullet_points(slide5, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4), bullets5, font_size=20)

# -------------------------------------------------------------
# SLIDE 6: TECH STACK
# -------------------------------------------------------------
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
setup_standard_slide(slide6, "TECH STACK")
bullets6 = [
    "🐍 Python: The core programming language used for the entire application.",
    "👁️ OpenCV: Powers camera input, image processing, and video feeds.",
    "✋ MediaPipe: Used for real-time hand tracking and 21 landmark detection.",
    "🧠 TensorFlow: The deep learning model handling gesture classification.",
    "🌐 Streamlit / Flask: Provides the web interface for user-friendly interaction.",
    "🔊 pyttsx3 / gTTS: The text-to-speech engine generating voice output."
]
add_bullet_points(slide6, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4), bullets6, font_size=20)

# -------------------------------------------------------------
# SLIDE 7: IMPLEMENTATION PLAN
# -------------------------------------------------------------
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
setup_standard_slide(slide7, "IMPLEMENTATION PLAN")
bullets7 = [
    "Phase 1: Collect and curate a dataset of sign language gestures.",
    "Phase 2: Use MediaPipe to extract 21 3D hand landmarks from each video frame.",
    "Phase 3: Train a machine learning classifier on landmark coordinates for gesture recognition.",
    "Phase 4: Integrate OpenCV for live camera feeds with the trained model for real-time prediction.",
    "Phase 5: Build a Streamlit or Flask web interface and add text-to-speech conversion.",
    "Phase 6: Conduct end-to-end testing, optimize performance, and prepare for a live demo."
]
add_bullet_points(slide7, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4), bullets7, font_size=18)

# -------------------------------------------------------------
# SLIDE 8: CONCLUSION
# -------------------------------------------------------------
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
setup_standard_slide(slide8, "CONCLUSION")
bullets8 = [
    "SignBridge successfully bridges the gap between sign language users and the hearing world through AI.",
    "Technology is most powerful when it makes the world more equitable and inclusive.",
    "Future enhancements will scale this impact by adding full sentence translation, multiple sign languages, and a dedicated mobile app."
]
add_bullet_points(slide8, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4), bullets8, font_size=24)

# -------------------------------------------------------------
# SLIDE 9: THANK YOU
# -------------------------------------------------------------
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_background_shape(slide9, prs, DARK_BG)
add_accent_bar(slide9, Inches(5.5), Inches(2.0), Inches(2.3), Inches(0.06), CORAL)

add_text_box(slide9, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.2),
             "THANK YOU", font_size=60, font_color=WHITE,
             bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide9, Inches(0), Inches(3.8), SLIDE_WIDTH, Inches(0.8),
             "SignBridge — AI Powered Sign Language Detector",
             font_size=22, font_color=CORAL, alignment=PP_ALIGN.CENTER)

add_text_box(slide9, Inches(0), Inches(4.5), SLIDE_WIDTH, Inches(1.5),
             "👥 Team: Reem Nisha & Team\n🏆 Hackathon 2025\n\nBuilt with Python | OpenCV | MediaPipe | TensorFlow | Streamlit\n🙋 Questions & Feedback Welcome!",
             font_size=18, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

add_accent_bar(slide9, Inches(0), Inches(6.8), SLIDE_WIDTH, Inches(0.7), TEAL)
add_text_box(slide9, Inches(0), Inches(6.95), SLIDE_WIDTH, Inches(0.55),
             "Breaking Barriers with AI | #SignBridge",
             font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# =============================================================
# SAVE PRESENTATION
# =============================================================
filename = "SignBridge_Final_Presentation.pptx"
prs.save(filename)
print(f"🎉 Presentation saved successfully as '{filename}'")